from flask import Flask, request, jsonify, send_file
from pptx import Presentation
import openai
import os
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from flask_jwt_extended import JWTManager, create_access_token, jwt_required, get_jwt_identity
from werkzeug.utils import secure_filename

app = Flask(__name__)

# Allow only GitHub Pages frontend to connect
CORS(app, resources={r"/*": {"origins": "https://ulagu2007.github.io"}})

# OpenAI API Key from environment variable (Replace in Render)
openai.api_key = os.getenv("OPENAI_API_KEY")

# Configurations
app.config["SQLALCHEMY_DATABASE_URI"] = "sqlite:///users.db"
app.config["JWT_SECRET_KEY"] = os.getenv("JWT_SECRET_KEY", "super-secret-key")  # Use env variable in Render
app.config["UPLOAD_FOLDER"] = "generated_ppts"

# Create Upload Folder
if not os.path.exists(app.config["UPLOAD_FOLDER"]):
    os.makedirs(app.config["UPLOAD_FOLDER"])

db = SQLAlchemy(app)
bcrypt = Bcrypt(app)
jwt = JWTManager(app)

# User Model
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(50), unique=True, nullable=False)
    password = db.Column(db.String(100), nullable=False)

# PPT History Model
class History(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    topic = db.Column(db.String(200), nullable=False)
    filename = db.Column(db.String(100), nullable=False)

# Generate AI PPT
def generate_ppt(topic):
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": f"Create a 5-slide PowerPoint presentation on {topic}. Provide only slide titles and bullet points."}]
    )

    slides = response["choices"][0]["message"]["content"].strip().split("\n\n")

    prs = Presentation()
    for slide_text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = slide_text.split("\n")
        title = lines[0].strip()
        content = "\n".join(lines[1:])

        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    ppt_filename = f"{topic.replace(' ', '_')}.pptx"
    ppt_path = os.path.join(app.config["UPLOAD_FOLDER"], ppt_filename)
    prs.save(ppt_path)

    return ppt_path

@app.route("/signup", methods=["POST"])
def signup():
    data = request.json
    username = data.get("username")
    password = bcrypt.generate_password_hash(data.get("password")).decode("utf-8")

    if User.query.filter_by(username=username).first():
        return jsonify({"error": "Username already exists"}), 400

    new_user = User(username=username, password=password)
    db.session.add(new_user)
    db.session.commit()
    return jsonify({"message": "Signup successful!"})

@app.route("/login", methods=["POST"])
def login():
    data = request.json
    user = User.query.filter_by(username=data.get("username")).first()

    if user and bcrypt.check_password_hash(user.password, data.get("password")):
        access_token = create_access_token(identity=user.id)
        return jsonify({"message": "Login successful!", "token": access_token})
    return jsonify({"error": "Invalid credentials"}), 401

@app.route("/generate", methods=["POST"])  # 🔥 Fixed API route
@jwt_required()
def generate_ppt_api():
    data = request.json
    topic = data.get("topic")

    if not topic:
        return jsonify({"error": "No topic provided"}), 400

    ppt_file = generate_ppt(topic)

    user_id = get_jwt_identity()
    history_entry = History(user_id=user_id, topic=topic, filename=ppt_file)
    db.session.add(history_entry)
    db.session.commit()

    return send_file(ppt_file, as_attachment=True)

@app.route("/history", methods=["GET"])
@jwt_required()
def get_history():
    user_id = get_jwt_identity()
    history = History.query.filter_by(user_id=user_id).all()
    return jsonify([{"topic": h.topic, "filename": h.filename} for h in history])

if __name__ == "__main__":
    app.run(debug=True)
